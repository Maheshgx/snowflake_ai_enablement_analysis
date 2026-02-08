#!/usr/bin/env python3
"""
Generate Snowflake AI Enablement Management Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths (relative to repository root)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.dirname(SCRIPT_DIR)
TEMPLATE_PATH = os.path.join(REPO_ROOT, "output/presentation/HMH Presentation Template_full.pptx")
OUTPUT_PATH = os.path.join(REPO_ROOT, "output/presentation/Snowflake_AI_Enablement_Management.pptx")

def add_notes(slide, notes_text):
    """Add speaker notes to a slide"""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text

def set_text_in_placeholder(slide, idx, text):
    """Set text in a placeholder by index"""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            shape.text = text
            return shape
    return None

def create_presentation():
    # Load template
    prs = Presentation(TEMPLATE_PATH)

    # Remove all existing slides (they're just examples)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # ========== SLIDE 1: Title Slide ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Intro Slide 1
    set_text_in_placeholder(slide, 12, "Snowflake AI Enablement")
    set_text_in_placeholder(slide, 13, "Strategic Opportunity Assessment\nFebruary 2026")
    add_notes(slide, """SPEAKER NOTES - Title Slide

Welcome and Introduction:
- Thank everyone for their time
- This presentation covers our AI enablement opportunity using Snowflake Cortex
- We've completed a comprehensive analysis of our data platform
- Key message: We have significant untapped AI potential in our existing data

Transition: Let's start with the agenda for today's discussion.""")

    # ========== SLIDE 2: Agenda ==========
    slide = prs.slides.add_slide(prs.slide_layouts[10])  # Agenda list
    set_text_in_placeholder(slide, 23, "Agenda")
    set_text_in_placeholder(slide, 10, "Executive Summary")
    set_text_in_placeholder(slide, 11, "AI Opportunity Assessment")
    set_text_in_placeholder(slide, 12, "Cost & Business Benefits")
    set_text_in_placeholder(slide, 13, "Implementation Roadmap")
    set_text_in_placeholder(slide, 14, "Next Steps & Recommendations")
    add_notes(slide, """SPEAKER NOTES - Agenda

Walk through the agenda:
1. Executive Summary - High-level findings and opportunity size
2. AI Opportunity Assessment - What we found in our data analysis
3. Cost & Business Benefits - The ROI case for AI enablement
4. Implementation Roadmap - Phased approach to capture value
5. Next Steps - What we're asking for today

This presentation is based on a comprehensive read-only analysis of our entire Snowflake environment.

Total time: ~20-25 minutes with Q&A""")

    # ========== SLIDE 3: Section Divider - Executive Summary ==========
    slide = prs.slides.add_slide(prs.slide_layouts[24])  # Section Divider - plain - blue
    set_text_in_placeholder(slide, 10, "Executive Summary")
    set_text_in_placeholder(slide, 11, "Key Findings & Opportunity Overview")
    add_notes(slide, """SPEAKER NOTES - Section Divider

Transition: Let's begin with the executive summary of our findings.""")

    # ========== SLIDE 4: The Opportunity ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "The AI Enablement Opportunity")
    set_text_in_placeholder(slide, 14, """We analyzed our entire Snowflake environment and identified massive AI potential:

• 149,342 AI-ready data candidates identified
• 67 databases with 424,000+ columns analyzed
• 5 Snowflake Cortex AI features applicable
• Zero additional infrastructure required
• SQL-based implementation - no ML expertise needed

Key Insight: Our data is already AI-ready. We just need to turn it on.""")
    add_notes(slide, """SPEAKER NOTES - The Opportunity

Key talking points:
- We conducted a comprehensive, read-only analysis of our entire Snowflake environment
- The analysis was automated and safe - only SELECT queries were run
- We found over 149,000 columns and tables that are ready for AI enablement TODAY

Emphasize:
- This is not about building new ML infrastructure
- Snowflake Cortex runs INSIDE Snowflake - no data movement
- Implementation is SQL-based - our existing team can do this
- The data is already there, already governed, already secured

This represents a significant competitive advantage we're not currently leveraging.""")

    # ========== SLIDE 5: Current State Metrics ==========
    slide = prs.slides.add_slide(prs.slide_layouts[43])  # Data - Icons - multicolor
    set_text_in_placeholder(slide, 26, "Current Data Landscape")
    set_text_in_placeholder(slide, 12, "67")
    set_text_in_placeholder(slide, 19, "937")
    set_text_in_placeholder(slide, 20, "13,927")
    set_text_in_placeholder(slide, 21, "424K+")
    set_text_in_placeholder(slide, 22, "Databases")
    set_text_in_placeholder(slide, 23, "Schemas")
    set_text_in_placeholder(slide, 24, "Tables/Views")
    set_text_in_placeholder(slide, 25, "Columns")
    add_notes(slide, """SPEAKER NOTES - Current State Metrics

Walk through the numbers:
- 67 databases - Fully analyzed with column-level data
- 937 schemas - Well-organized data domains
- 13,927 tables and views - Rich data assets
- 424,000+ columns - Massive analytical potential

Context:
- This represents the fully analyzed portion of our data estate
- We've invested significantly in getting data into Snowflake
- Now it's time to extract more value from that investment
- AI enablement is the next logical step in our data journey

The question isn't IF we should enable AI, but HOW FAST we can do it.""")

    # ========== SLIDE 6: AI Candidates by Feature ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "AI Opportunities by Cortex Feature")
    set_text_in_placeholder(slide, 14, """Cortex LLM Functions: 125,468 candidates
• Text summarization, classification, sentiment analysis
• Immediate value for content and feedback processing

Cortex Search / RAG: 9,925 candidates
• Semantic search across educational content
• Natural language queries on documentation

Cortex ML: 6,216 candidates
• Time-series forecasting and anomaly detection
• Predictive analytics on operational data

Cortex Extract: 2,566 candidates
• Semi-structured data processing (JSON, XML)

Document AI: 269+ stages
• PDF and document processing from cloud storage""")
    add_notes(slide, """SPEAKER NOTES - AI Candidates by Feature

Explain each Cortex feature:

1. Cortex LLM (125K candidates) - HIGHEST PRIORITY
   - These are text columns we can immediately enrich
   - Summarize long content, classify topics, analyze sentiment
   - Use cases: Auto-tag content, analyze student feedback, summarize reports

2. Cortex Search (10K candidates)
   - Enable natural language search across our content
   - "Find algebra word problems for 8th grade" instead of keyword search
   - Huge improvement for content discovery

3. Cortex ML (6K candidates)
   - Tables with timestamps and metrics ready for forecasting
   - Predict enrollment trends, detect anomalies in system metrics
   - Built-in models, no data science team needed

4. Cortex Extract (2.5K candidates)
   - JSON/XML columns we can parse and analyze
   - Unlock value from semi-structured data

5. Document AI (269 stages)
   - We have 269 cloud storage stages with documents
   - Can extract text and insights from PDFs, Word docs, images""")

    # ========== SLIDE 7: Section Divider - Benefits ==========
    slide = prs.slides.add_slide(prs.slide_layouts[25])  # Section Divider - plain - green
    set_text_in_placeholder(slide, 10, "Cost & Business Benefits")
    set_text_in_placeholder(slide, 11, "ROI Analysis & Value Drivers")
    add_notes(slide, """SPEAKER NOTES - Section Divider

Transition: Now let's look at the financial and business case for AI enablement.""")

    # ========== SLIDE 8: Cost Benefits ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "Cost Reduction Opportunities")
    set_text_in_placeholder(slide, 14, """Infrastructure Savings: 40-60% reduction
• No separate ML platforms or GPU clusters needed
• Cortex runs natively inside Snowflake

Data Movement Elimination: 70-80% reduction
• Process data where it lives - zero ETL to external tools
• No data copying, no sync issues, no additional storage

Development Acceleration: 60-80% faster
• SQL-based AI vs. months of ML engineering
• Existing team can implement immediately

Tool Consolidation: 30-50% reduction
• Replace multiple AI/ML tool licenses
• Single platform for data + AI""")
    add_notes(slide, """SPEAKER NOTES - Cost Benefits

Infrastructure (40-60% savings):
- Traditional ML requires separate compute clusters, often with expensive GPUs
- Cortex uses Snowflake's existing compute - we already pay for this
- No new infrastructure to provision, secure, or maintain

Data Movement (70-80% savings):
- Moving data to external ML platforms is expensive and risky
- ETL pipelines, storage duplication, sync jobs
- With Cortex, AI runs WHERE the data already is
- This also reduces security risk - data never leaves Snowflake

Development Time (60-80% faster):
- Traditional ML: data prep, model training, deployment, monitoring
- Cortex: Write a SQL query, get AI results
- Our existing SQL developers can do this TODAY

Tool Consolidation (30-50% savings):
- We likely have multiple point solutions for various AI tasks
- Cortex consolidates these into one platform
- Single vendor, single skill set, single security model""")

    # ========== SLIDE 9: Business Value ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "Business Value & Impact")
    set_text_in_placeholder(slide, 14, """Content Auto-Classification
• 90% reduction in manual content tagging effort
• Consistent, scalable categorization

Semantic Search Enhancement
• 3x improvement in content discovery
• Natural language queries for users

Predictive Analytics
• 20-30% better forecasting accuracy
• Data-driven planning and decisions

Document Processing Automation
• 80% reduction in manual document review
• Extract insights from unstructured content

Anomaly Detection
• 50% faster issue identification
• Proactive problem resolution""")
    add_notes(slide, """SPEAKER NOTES - Business Value

For each benefit, provide concrete examples:

Content Auto-Classification (90% manual reduction):
- Currently: Manual tagging of educational content
- With AI: Automatic classification by subject, grade, difficulty
- Impact: Free up content team for higher-value work

Semantic Search (3x improvement):
- Currently: Keyword search, often misses relevant content
- With AI: "Find content about fractions for struggling 5th graders"
- Impact: Teachers find what they need faster, better outcomes

Predictive Analytics (20-30% better):
- Currently: Spreadsheet-based forecasting
- With AI: ML-powered predictions with confidence intervals
- Impact: Better resource planning, reduced surprises

Document Processing (80% reduction):
- Currently: Manual review of uploaded documents
- With AI: Automatic extraction and categorization
- Impact: Faster processing, consistent quality

Anomaly Detection (50% faster):
- Currently: Reactive - find issues when users complain
- With AI: Proactive alerts on unusual patterns
- Impact: Fix problems before they impact users""")

    # ========== SLIDE 10: ROI Summary ==========
    slide = prs.slides.add_slide(prs.slide_layouts[43])  # Data - Icons - multicolor
    set_text_in_placeholder(slide, 26, "Expected Return on Investment")
    set_text_in_placeholder(slide, 12, "150-400%")
    set_text_in_placeholder(slide, 19, "3-6")
    set_text_in_placeholder(slide, 20, "$100-500K")
    set_text_in_placeholder(slide, 21, "90%+")
    set_text_in_placeholder(slide, 22, "Year 1 ROI")
    set_text_in_placeholder(slide, 23, "Month Payback")
    set_text_in_placeholder(slide, 24, "Annual Savings")
    set_text_in_placeholder(slide, 25, "Risk Reduction")
    add_notes(slide, """SPEAKER NOTES - ROI Summary

Key metrics to emphasize:

150-400% Year 1 ROI:
- Conservative estimate based on industry benchmarks
- Actual ROI depends on which use cases we prioritize
- Higher ROI comes from high-volume, repetitive tasks

3-6 Month Payback:
- Quick wins can show value in weeks
- Full payback typically within two quarters
- Low risk given minimal upfront investment

$100-500K Annual Savings:
- Range depends on current manual effort being replaced
- Includes: tool consolidation, labor efficiency, faster time-to-insight
- Does NOT include revenue opportunities from better analytics

90%+ Risk Reduction:
- Data stays in Snowflake - no new attack surface
- Existing governance and security controls apply
- Audit trail for all AI operations

This is a LOW-RISK, HIGH-REWARD initiative.""")

    # ========== SLIDE 11: Section Divider - Roadmap ==========
    slide = prs.slides.add_slide(prs.slide_layouts[22])  # Section Divider - plain - yellow
    set_text_in_placeholder(slide, 10, "Implementation Roadmap")
    set_text_in_placeholder(slide, 11, "Phased Approach to Value")
    add_notes(slide, """SPEAKER NOTES - Section Divider

Transition: Let me walk you through our recommended implementation approach.""")

    # ========== SLIDE 12: Three-Phase Approach ==========
    slide = prs.slides.add_slide(prs.slide_layouts[41])  # Timeline / Progression
    set_text_in_placeholder(slide, 26, "Three-Phase Implementation")
    set_text_in_placeholder(slide, 16, "Phase 1")
    set_text_in_placeholder(slide, 17, "Phase 2")
    set_text_in_placeholder(slide, 18, "Phase 3")
    set_text_in_placeholder(slide, 19, "")
    set_text_in_placeholder(slide, 20, "")
    set_text_in_placeholder(slide, 21, "Quick Wins\n1-2 Weeks")
    set_text_in_placeholder(slide, 22, "Core Features\n2-4 Weeks")
    set_text_in_placeholder(slide, 23, "Advanced\n1-2 Months")
    set_text_in_placeholder(slide, 24, "")
    set_text_in_placeholder(slide, 25, "")
    add_notes(slide, """SPEAKER NOTES - Three-Phase Approach

Phase 1: Quick Wins (1-2 Weeks)
- Immediate value, minimal risk
- Text summarization on marketing content
- Sentiment analysis on feedback data
- Proof of concept for stakeholders
- Estimated: 40 implementation hours, 500 credits/month

Phase 2: Core Features (2-4 Weeks)
- Build on Phase 1 success
- Cortex Search for content discovery
- Classification for content tagging
- Basic forecasting models
- Estimated: 120 implementation hours, 2,000 credits/month

Phase 3: Advanced Capabilities (1-2 Months)
- Full-scale deployment
- Document AI on cloud stages
- Anomaly detection on operational data
- Complex ML pipelines
- Estimated: 200 implementation hours, 5,000 credits/month

Key point: Each phase delivers standalone value. We can pause and assess after any phase.""")

    # ========== SLIDE 13: Phase 1 Details ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "Phase 1: Quick Wins (1-2 Weeks)")
    set_text_in_placeholder(slide, 14, """Target Use Cases:
• Content summarization on marketing materials
• Sentiment analysis on customer feedback
• Basic text classification for content tagging

Why Start Here:
• Highest-scoring candidates from our analysis
• Minimal governance complexity
• Immediate, visible results
• Low credit consumption

Success Metrics:
• Content processing time reduced by 80%
• Query latency under 2 seconds
• Stakeholder validation of AI quality""")
    add_notes(slide, """SPEAKER NOTES - Phase 1 Details

Specific targets for Phase 1:
- Marketing content in ANNUITAS database - 446 avg character length, perfect for summarization
- Customer feedback tables - sentiment scoring
- Course content - classification by subject area

Why these are "quick wins":
1. Data is clean and well-structured
2. No PII concerns - marketing and content data
3. Clear before/after comparison
4. Results are easy to validate

Resource requirements:
- 40 hours of implementation time
- XS warehouse sufficient
- ~500 credits/month ongoing
- No additional headcount

Success criteria:
- Demonstrate 80%+ time savings on specific tasks
- Sub-2-second response times
- Business users validate quality
- Document lessons learned for Phase 2""")

    # ========== SLIDE 14: Governance Considerations ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "Governance & Risk Management")
    set_text_in_placeholder(slide, 14, """Data Protection:
• 7,784 high-risk columns identified (PII/sensitive)
• AI enablement requires governance review first
• Existing Snowflake RBAC and masking policies apply

Compliance Considerations:
• FERPA - Student data protection
• COPPA - Children's privacy
• GDPR - European data subjects

Recommended Safeguards:
• Dynamic data masking before AI processing
• Row-level security for sensitive tables
• Cortex Guard for content filtering
• Complete audit trail of all AI operations""")
    add_notes(slide, """SPEAKER NOTES - Governance

This is critical - we're not ignoring governance:

PII Identification:
- Our analysis flagged 7,784 columns with potential PII
- Includes: email, names, addresses, phone numbers, birth dates
- These require extra review before AI enablement
- Phase 1 intentionally avoids these columns

Compliance requirements:
- FERPA: Student education records are protected
- COPPA: Extra protections for children under 13
- GDPR: EU data subject rights
- Our existing compliance framework extends to AI

Built-in protections:
- Snowflake's RBAC works with Cortex - same access controls
- Dynamic masking can anonymize data before AI processing
- Row-level security limits what AI can "see"
- Every AI query is logged for audit

Key message: We're not asking to bypass governance. We're proposing to extend our existing governance to AI workloads.""")

    # ========== SLIDE 15: Section Divider - Next Steps ==========
    slide = prs.slides.add_slide(prs.slide_layouts[26])  # Section Divider - plain - purple
    set_text_in_placeholder(slide, 10, "Next Steps")
    set_text_in_placeholder(slide, 11, "Recommendations & Ask")
    add_notes(slide, """SPEAKER NOTES - Section Divider

Transition: Let me conclude with our recommendations and what we're asking for today.""")

    # ========== SLIDE 16: Recommendations ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "Our Recommendations")
    set_text_in_placeholder(slide, 14, """1. Approve Phase 1 Quick Wins
   • 2-week pilot with defined success criteria
   • Minimal investment: 40 hours + 500 credits/month

2. Establish AI Governance Framework
   • Extend existing data governance to AI workloads
   • Define approval process for new AI use cases

3. Identify Business Champions
   • Each department nominates AI use case owner
   • Build internal expertise and advocacy

4. Schedule Phase 1 Review
   • Two weeks post-implementation
   • Go/no-go decision for Phase 2""")
    add_notes(slide, """SPEAKER NOTES - Recommendations

Be specific about what you're asking for:

1. Phase 1 Approval:
   - This is a low-risk pilot
   - 40 hours = roughly 1 week of one developer's time
   - 500 credits/month = approximately $X (check current pricing)
   - Clear success criteria defined upfront

2. Governance Framework:
   - We're not asking to create something new
   - Extend existing data governance policies
   - Add AI-specific considerations to review process
   - Security team involvement from the start

3. Business Champions:
   - AI succeeds when business owns the use cases
   - Each department should identify their top opportunity
   - Creates accountability and ensures relevance
   - Builds internal advocates

4. Review Checkpoint:
   - Built-in decision point
   - If Phase 1 doesn't deliver, we pause
   - Data-driven go/no-go for continued investment
   - De-risks the overall initiative""")

    # ========== SLIDE 17: The Ask ==========
    slide = prs.slides.add_slide(prs.slide_layouts[30])  # Content - 1 column bullets
    set_text_in_placeholder(slide, 13, "What We Need Today")
    set_text_in_placeholder(slide, 14, """Approval to Proceed:
• Green light for Phase 1 implementation
• Budget allocation: ~$5,000 for pilot period

Resource Commitment:
• 1 developer for 2 weeks (40 hours)
• 1 business analyst for use case validation
• Security/governance review participation

Timeline:
• Week 1-2: Phase 1 implementation
• Week 3: Results review and Phase 2 planning
• Week 4+: Phase 2 execution (pending approval)

Decision Needed: Approve Phase 1 pilot program""")
    add_notes(slide, """SPEAKER NOTES - The Ask

Be direct about what you need:

Budget:
- Approximately $5,000 for the pilot period
- Covers: Snowflake credits + implementation time
- This is a rounding error compared to potential savings

Resources:
- 1 developer - can be existing team member
- 1 business analyst - to validate results make sense
- Security participation - review, not full-time

Timeline commitment:
- 2 weeks to prove value
- Decision checkpoint at week 3
- No long-term commitment required upfront

What you're NOT asking for:
- New headcount
- New tools or platforms
- Large capital investment
- Multi-year commitment

This is a low-risk, time-boxed pilot with clear success criteria.

PAUSE FOR QUESTIONS/DISCUSSION before moving to Q&A slide.""")

    # ========== SLIDE 18: Questions ==========
    slide = prs.slides.add_slide(prs.slide_layouts[54])  # Questions - no image
    set_text_in_placeholder(slide, 10, "Questions?")
    add_notes(slide, """SPEAKER NOTES - Q&A

Anticipated questions and answers:

Q: How does this affect our Snowflake costs?
A: Cortex uses credits like any compute. We've estimated 500-5,000 credits/month depending on phase. This is included in our budget ask.

Q: What about data security?
A: Data never leaves Snowflake. All existing security controls apply. Every AI operation is logged. We're actually MORE secure than external AI tools.

Q: Do we need to hire ML engineers?
A: No. Cortex is SQL-based. Our existing developers can implement this. That's one of the key benefits.

Q: What if the AI gives wrong answers?
A: Phase 1 focuses on low-risk use cases. Human review is part of the process. We're augmenting, not replacing, human judgment.

Q: How does this compare to ChatGPT/OpenAI?
A: Cortex runs INSIDE our data platform. No data sent externally. Enterprise-grade security. Purpose-built for our data.

Q: What's the risk if we don't do this?
A: Competitors are already using AI. Our data is a strategic asset that's underutilized. The cost of inaction is falling behind.

THANK THEM FOR THEIR TIME""")

    # ========== SLIDE 19: Closing ==========
    slide = prs.slides.add_slide(prs.slide_layouts[63])  # HMH closing slide
    add_notes(slide, """SPEAKER NOTES - Closing

Final remarks:
- Thank everyone for their time and attention
- Reiterate the key message: Low risk, high reward opportunity
- Our data is already AI-ready - we just need to activate it
- Request a decision on Phase 1 approval

Follow-up:
- Offer to send the presentation and supporting documentation
- Schedule follow-up meeting if needed
- Provide contact information for questions

Leave-behind materials:
- Executive Summary document
- AI Strategy Roadmap
- Detailed analysis reports (available on request)""")

    # Save presentation
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
